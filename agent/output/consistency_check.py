# agent/output/consistency_check.py — Vérifie la cohérence Excel ↔ Présentation
import openpyxl
from pptx import Presentation


def verify_consistency(excel_path: str, pptx_path: str) -> dict:
    """
    Vérifie que les chiffres de la présentation correspondent au classeur Excel.
    
    Compare tous les nombres trouvés dans les deux fichiers.
    
    Returns:
        Dictionnaire avec le statut et les incohérences détectées.
    """
    # Charger le classeur Excel
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    
    # Extraire tous les nombres du Excel
    excel_numbers = set()
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, (int, float)):
                    excel_numbers.add(round(float(cell.value), 2))
    
    # Charger la présentation
    prs = Presentation(pptx_path)
    
    # Extraire tous les nombres de la présentation
    pptx_numbers = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Remplacer virgule par point et enlever les %
                        text = run.text.replace(",", ".").replace("%", "").strip()
                        try:
                            num = float(text)
                            pptx_numbers.add(round(num, 2))
                        except ValueError:
                            pass  # Ce n'est pas un nombre, on ignore
            # Vérifier aussi les données des graphiques
            if shape.has_chart:
                for series in shape.chart.series:
                    for val in series.values:
                        if val is not None:
                            pptx_numbers.add(round(float(val), 2))
    
    # Identifier les incohérences (nombres dans la présentation mais pas dans Excel)
    pptx_only = pptx_numbers - excel_numbers
    # Filtrer les zéros et les nombres triviaux
    inconsistencies = [n for n in pptx_only if abs(n) > 0.01]
    
    return {
        "passed": len(inconsistencies) == 0,
        "excel_count": len(excel_numbers),
        "pptx_count": len(pptx_numbers),
        "inconsistencies": inconsistencies[:10],
    }
