# agent/output/pptx_generator.py — Génération de présentations PowerPoint
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


class PresentationGenerator:
    """
    Génère une présentation de 6 slides à partir des résultats de l'agent.

    Slides :
    1. Titre et problème
    2. Résultats clés
    3. Validation
    4. Recommandations (résumé)
    5. Détail des recommandations
    6. Remerciement
    """
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
    
    def _add_title_slide(self, title: str, subtitle: str):
        """Slide 1 : Titre et problème."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Layout vide
        
        # Titre
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x2F, 0x54, 0x96)
        
        # Sous-titre
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
    
    def _add_findings_slide(self, title: str, findings: list[str]):
        """Slide avec une liste de points."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Titre
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        
        # Points
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        
        for i, finding in enumerate(findings):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = f"• {finding}"
            p.font.size = Pt(18)
            p.space_after = Pt(12)
    
    def generate(
        self,
        problem: str,
        findings: list[str],
        validation: dict,
        recommendations: list[dict],
        output_path: str,
        title: str = "Rapport d'analyse",
    ) -> str:
        """Génère la présentation complète de 6 slides."""

        # Slide 1 : Titre
        self._add_title_slide(title, problem)
        
        # Slide 2 : Résultats clés
        self._add_findings_slide("Résultats Clés", findings)
        
        # Slide 3 : Validation
        validation_points = [
            f"{k}: {'OK' if v.get('passed') else 'ÉCHEC'} — {v.get('result', '')}"
            for k, v in validation.items()
        ]
        self._add_findings_slide("Validation", validation_points)
        
        # Slide 4 : Recommandations
        rec_points = [
            f"{rec.get('title', '')} (Impact: {rec.get('impact', 'N/A')}, "
            f"Délai: {rec.get('timeline', 'N/A')})"
            for rec in recommendations
        ]
        self._add_findings_slide("Recommandations", rec_points)
        
        # Slide 5 : Détail des recommandations
        detail_points = [rec.get("description", "") for rec in recommendations]
        self._add_findings_slide("Détail des Recommandations", detail_points)
        
        # Slide 6 : Remerciement
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Merci — Questions ?"
        p.font.size = Pt(36)
        p.font.bold = True
        
        self.prs.save(output_path)
        return output_path
