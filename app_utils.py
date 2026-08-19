# app_utils.py : fonctions pures de l'application Streamlit (nettoyage,
# exports, sécurité), séparées de app.py pour être testables sans
# dépendre du runtime Streamlit (app.py exécute du code au niveau module).
import csv
import io
import os
import re
import zipfile

import numpy as np
import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Font
from openpyxl.utils.dataframe import dataframe_to_rows
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from agent.tools.spreadsheet_safety import neutralize_formulas, unique_sheet_title

DATE_KEYWORDS = ["date", "time", "jour", "mois", "annee", "année", "year"]
DECIMAL_PATTERN = re.compile(r"^-?\d+,\d+$")
LEADING_ZERO_PATTERN = re.compile(r"^0\d")


def read_csv_robust(uploaded_file) -> tuple[pd.DataFrame, str]:
    """Lit un CSV en devinant l'encodage et le séparateur (utile pour les exports Excel FR).

    Le séparateur est deviné avec un jeu de délimiteurs restreint (évite que le
    détecteur automatique de pandas choisisse un caractère au hasard sur un
    fichier à une seule colonne) ; à défaut de détection fiable, on retombe sur
    la virgule.

    Toutes les colonnes sont lues en texte brut (`dtype=str`) : sans ça, le
    parseur C de pandas infère lui-même les colonnes numériques et tronque les
    zéros non significatifs ("001" -> 1) avant même que `clean_data` ait la
    main pour les protéger. `clean_data` se charge ensuite de reconvertir en
    nombre/date les colonnes qui le méritent réellement.
    """
    last_err = None
    for enc in ["utf-8-sig", "utf-8", "latin-1"]:
        try:
            uploaded_file.seek(0)
            sample_bytes = uploaded_file.read(8192)
            sample_text = sample_bytes.decode(enc, errors="ignore")
            uploaded_file.seek(0)
            try:
                sep = csv.Sniffer().sniff(sample_text, delimiters=",;\t|").delimiter
            except csv.Error:
                sep = ","
            df = pd.read_csv(uploaded_file, sep=sep, encoding=enc, dtype=str)
            return df, enc
        except Exception as e:
            last_err = e
            continue
    raise last_err


def clean_data(df: pd.DataFrame) -> tuple[pd.DataFrame, dict]:
    """Nettoie un DataFrame : types, doublons, texte, dates, décimales FR."""
    df = df.copy()
    original_rows, original_cols = df.shape

    # Noms de colonnes propres
    df.columns = [str(c).strip() for c in df.columns]

    # Lignes/colonnes entièrement vides
    df = df.dropna(axis=0, how="all")
    df = df.dropna(axis=1, how="all")
    rows_after_empty_drop = df.shape[0]

    # Texte : strip + chaînes vides -> NaN
    obj_cols = df.select_dtypes(include=["object", "str"]).columns
    for col in obj_cols:
        df[col] = df[col].apply(lambda x: x.strip() if isinstance(x, str) else x)
        df[col] = df[col].replace({"": np.nan})

    # Nombres au format européen ("120,50" -> 120.50)
    converted_decimal = []
    for col in df.select_dtypes(include=["object", "str"]).columns:
        sample = df[col].dropna().astype(str)
        if len(sample) > 0 and sample.str.match(DECIMAL_PATTERN).mean() > 0.7:
            df[col] = df[col].astype(str).str.replace(",", ".", regex=False)
            converted_decimal.append(col)

    # Colonnes numériques stockées en texte (on protège les identifiants à
    # zéros non significatifs, ex. codes postaux "00501", pour ne pas les
    # convertir en nombres et perdre les zéros de tête)
    converted_numeric = []
    for col in df.select_dtypes(include=["object", "str"]).columns:
        sample = df[col].dropna().astype(str)
        if len(sample) > 0 and sample.str.match(LEADING_ZERO_PATTERN).mean() > 0.3:
            continue
        converted = pd.to_numeric(df[col], errors="coerce")
        non_null = df[col].notna().sum()
        if non_null > 0 and converted.notna().sum() / non_null > 0.9:
            df[col] = converted
            converted_numeric.append(col)

    # Colonnes de dates (formats mixtes ISO / JJ-MM / MM-JJ)
    converted_dates = []
    for col in df.select_dtypes(include=["object", "str"]).columns:
        if any(k in col.lower() for k in DATE_KEYWORDS):
            non_null = df[col].notna().sum()
            if non_null == 0:
                continue
            # ISO (AAAA-MM-JJ) est essayé en premier et isolément : avec
            # format="mixed", dayfirst=True inverse à tort jour/mois même sur
            # des dates ISO non ambiguës dès que le jour est <= 12 (bug pandas
            # observé). ISO8601 lève cette ambiguïté puisque l'année est
            # toujours le premier composant.
            parsed = pd.to_datetime(df[col], errors="coerce", format="ISO8601")
            remaining = parsed.isna() & df[col].notna()
            if remaining.any():
                # Le reste (formats JJ/MM/AAAA ou MM/JJ/AAAA) reste ambigu :
                # on garde l'heuristique par comparaison des deux lectures.
                parsed_eu = pd.to_datetime(df[col][remaining], errors="coerce", format="mixed", dayfirst=True)
                parsed_us = pd.to_datetime(df[col][remaining], errors="coerce", format="mixed", dayfirst=False)
                fallback = parsed_eu if parsed_eu.notna().sum() >= parsed_us.notna().sum() else parsed_us
                parsed.loc[remaining] = fallback
            if parsed.notna().sum() / non_null > 0.7:
                df[col] = parsed
                converted_dates.append(col)

    # Doublons exacts
    duplicates_removed = int(df.duplicated().sum())
    df = df.drop_duplicates().reset_index(drop=True)

    missing_after = df.isna().sum()
    missing_after = {k: int(v) for k, v in missing_after[missing_after > 0].items()}

    report = {
        "original_rows": original_rows,
        "original_cols": original_cols,
        "final_rows": df.shape[0],
        "final_cols": df.shape[1],
        "empty_rows_removed": original_rows - rows_after_empty_drop,
        "duplicates_removed": duplicates_removed,
        "columns_converted_decimal": converted_decimal,
        "columns_converted_numeric": converted_numeric,
        "columns_converted_date": converted_dates,
        "missing_values": missing_after,
    }
    return df, report


def sanitize_sheet_name(name: str, used: set) -> str:
    """Dérive un nom d'onglet Excel valide et unique à partir d'un nom de fichier.

    Retire l'extension puis délègue à unique_sheet_title() pour la partie
    commune avec le mode agent (troncature à 31 caractères, caractères
    interdits, dédoublonnage).
    """
    return unique_sheet_title(name.rsplit(".", 1)[0], used)


def build_markdown_report(query: str, model_used: str, files: list, analysis_text: str, timestamp: str) -> str:
    """Assemble un rapport Markdown autonome (question + données + analyse)."""
    files_desc = "\n".join(f"- {name} : {len(df)} lignes" for name, df in files)
    return "\n".join([
        f"# Rapport d'analyse : {timestamp}",
        "",
        f"**Question :** {query}",
        f"**Modèle utilisé :** {model_used}",
        "**Fichiers analysés (après nettoyage) :**",
        files_desc,
        "",
        "## Analyse",
        "",
        analysis_text,
    ])


def build_excel_report(files: list, describe_list: list, query: str, model_used: str, analysis_text: str) -> bytes:
    """Génère un classeur Excel en mémoire : pour chaque fichier nettoyé, sa
    feuille de données suivie d'une feuille de statistiques descriptives
    (le même `df.describe()` que celui déjà affiché à l'écran, pas
    recalculé ici, pour garantir que le fichier téléchargé montre
    exactement ce que l'utilisateur a déjà vu), puis une feuille d'analyse.
    """
    wb = Workbook()
    wb.remove(wb.active)
    used_names = set()

    for (name, df), desc in zip(files, describe_list):
        ws = wb.create_sheet(sanitize_sheet_name(name, used_names))
        for row in dataframe_to_rows(neutralize_formulas(df), index=False, header=True):
            ws.append(row)
        for cell in ws[1]:
            cell.font = Font(bold=True)

        stem = name.rsplit(".", 1)[0]
        ws_stats = wb.create_sheet(unique_sheet_title(f"Stats - {stem}", used_names))
        # openpyxl.dataframe_to_rows(index=True, header=True) insère toujours
        # une ligne fantôme juste après l'en-tête (le nom de l'index sur sa
        # propre ligne), un artefact de mise en forme, pas une donnée à garder.
        for i, row in enumerate(dataframe_to_rows(neutralize_formulas(desc), index=True, header=True)):
            if i != 1:
                ws_stats.append(row)
        for cell in ws_stats[1]:
            cell.font = Font(bold=True)

    ws_analysis = wb.create_sheet("Analyse")
    ws_analysis.column_dimensions["A"].width = 100
    ws_analysis["A1"] = "Question"
    ws_analysis["A1"].font = Font(bold=True)
    ws_analysis["A2"] = query
    ws_analysis["A4"] = "Modèle"
    ws_analysis["A4"].font = Font(bold=True)
    ws_analysis["A5"] = model_used
    ws_analysis["A7"] = "Analyse"
    ws_analysis["A7"].font = Font(bold=True)
    for i, line in enumerate(analysis_text.split("\n")):
        ws_analysis.cell(row=8 + i, column=1, value=line)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def build_csv_export(files: list) -> tuple[bytes, str, str]:
    """Exporte les données nettoyées : un CSV seul, ou un ZIP si plusieurs fichiers."""
    if len(files) == 1:
        name, df = files[0]
        return neutralize_formulas(df).to_csv(index=False).encode("utf-8-sig"), "donnees_nettoyees.csv", "text/csv"

    buf = io.BytesIO()
    used_names = set()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for name, df in files:
            # os.path.basename() écarte tout composant de chemin qu'un nom de
            # fichier téléversé pourrait contenir (ex. "../../evil.csv") :
            # l'entrée ZIP ne doit jamais pouvoir pointer hors du dossier
            # d'extraction attendu par qui l'ouvrira.
            base = re.sub(r"\.csv$", "", os.path.basename(name), flags=re.IGNORECASE)
            csv_name = f"{base}_nettoye.csv"
            i = 2
            while csv_name.lower() in used_names:
                csv_name = f"{base}_nettoye_{i}.csv"
                i += 1
            used_names.add(csv_name.lower())
            zf.writestr(csv_name, neutralize_formulas(df).to_csv(index=False))
    return buf.getvalue(), "donnees_nettoyees.zip", "application/zip"


def _chunk_text(text: str, max_chars: int = 700) -> list:
    """Découpe un texte en morceaux d'environ max_chars pour tenir sur des diapositives lisibles."""
    paragraphs = [p for p in text.split("\n") if p.strip()]
    chunks, current, current_len = [], [], 0
    for p in paragraphs:
        if current_len + len(p) > max_chars and current:
            chunks.append("\n".join(current))
            current, current_len = [], 0
        current.append(p)
        current_len += len(p)
    if current:
        chunks.append("\n".join(current))
    return chunks or [""]


def build_pptx_report(query: str, model_used: str, files: list, analysis_text: str, timestamp: str) -> bytes:
    """Génère une présentation PowerPoint générique : titre, fichiers analysés, puis l'analyse découpée en diapositives."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.3), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Rapport d'analyse"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x36, 0x25, 0x5C)

    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(2.5))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = query
    tf2.paragraphs[0].font.size = Pt(20)
    p3 = tf2.add_paragraph()
    p3.text = f"{timestamp} · {model_used}"
    p3.font.size = Pt(14)

    slide2 = prs.slides.add_slide(blank)
    box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = "Fichiers analysés"
    p.font.size = Pt(32)
    p.font.bold = True
    box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    for i, (name, df) in enumerate(files):
        para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        para.text = f"- {name} : {len(df)} lignes"
        para.font.size = Pt(18)

    chunks = _chunk_text(analysis_text)
    for i, chunk in enumerate(chunks):
        slide3 = prs.slides.add_slide(blank)
        title = "Analyse" if len(chunks) == 1 else f"Analyse ({i + 1}/{len(chunks)})"
        tbox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tp = tbox.text_frame.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(28)
        tp.font.bold = True

        cbox = slide3.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.9))
        tf2 = cbox.text_frame
        tf2.word_wrap = True
        for j, line in enumerate(chunk.split("\n")):
            para = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
            para.text = line
            para.font.size = Pt(14)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
