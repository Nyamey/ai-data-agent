# tests/test_pptx_generator.py : génération de présentations PowerPoint
#
# Déjà exercé indirectement par tests/test_consistency_check.py, mais avec
# des findings/validation/recommendations presque vides -- la plupart du
# contenu des diapositives (boucles sur plusieurs éléments) restait non
# couverte. Ici avec des listes réellement peuplées (plusieurs findings,
# plusieurs recommandations, un mélange de contrôles OK/ÉCHEC).
from pptx import Presentation

from agent.output.pptx_generator import PresentationGenerator


def _generate(tmp_path, **overrides):
    kwargs = dict(
        problem="Pourquoi la rétention baisse-t-elle ?",
        findings=["Rétention en baisse de 15% sur mobile", "Stable sur desktop", "Web en légère hausse"],
        validation={
            "doublons": {"result": 0, "passed": True},
            "plage_dates": {"result": "2024-01-01 -> 2024-06-30", "passed": True},
            "coherence": {"result": "écart détecté", "passed": False},
        },
        recommendations=[
            {"title": "Améliorer l'onboarding mobile", "description": "Détail 1.",
             "impact": "élevé", "feasibility": "moyenne", "timeline": "court terme"},
            {"title": "Revoir les notifications push", "description": "Détail 2.",
             "impact": "moyen", "feasibility": "élevée", "timeline": "moyen terme"},
        ],
        output_path=str(tmp_path / "presentation.pptx"),
    )
    kwargs.update(overrides)
    path = PresentationGenerator().generate(**kwargs)
    return path, Presentation(path)


def _all_text(slide):
    return "\n".join(
        shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
    )


def test_generate_produces_six_slides(tmp_path):
    _, prs = _generate(tmp_path)
    assert len(prs.slides) == 6


def test_title_slide_shows_problem_and_custom_title(tmp_path):
    _, prs = _generate(tmp_path, title="Rapport personnalisé")
    text = _all_text(prs.slides[0])
    assert "Rapport personnalisé" in text
    assert "Pourquoi la rétention baisse-t-elle ?" in text


def test_findings_slide_lists_every_finding(tmp_path):
    _, prs = _generate(tmp_path)
    text = _all_text(prs.slides[1])
    assert "Rétention en baisse de 15% sur mobile" in text
    assert "Stable sur desktop" in text
    assert "Web en légère hausse" in text


def test_validation_slide_shows_pass_and_fail_markers(tmp_path):
    _, prs = _generate(tmp_path)
    text = _all_text(prs.slides[2])
    assert "OK" in text
    assert "ÉCHEC" in text
    assert "doublons" in text
    assert "coherence" in text


def test_recommendations_summary_slide_lists_all_titles(tmp_path):
    _, prs = _generate(tmp_path)
    text = _all_text(prs.slides[3])
    assert "Améliorer l'onboarding mobile" in text
    assert "Revoir les notifications push" in text
    assert "élevé" in text


def test_recommendations_detail_slide_shows_descriptions(tmp_path):
    _, prs = _generate(tmp_path)
    text = _all_text(prs.slides[4])
    assert "Détail 1." in text
    assert "Détail 2." in text


def test_generate_with_empty_lists_does_not_crash(tmp_path):
    path, prs = _generate(tmp_path, findings=[], validation={}, recommendations=[])
    assert len(prs.slides) == 6
