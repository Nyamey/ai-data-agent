# tests/test_export_node.py — Nœud 8 : export Excel/PowerPoint
#
# export_node n'avait jamais été testé directement (0 % de couverture avant
# ce fichier) : il n'était exercé que par des vérifications manuelles hors
# suite pytest. Les états ici sont construits en chaînant les vrais nœuds
# (build/test/validate) sur une vraie base DuckDB, pour que weekly_retention/
# driver_analysis contiennent de vraies requêtes "query" -- export_node les
# ré-exécute pour obtenir des DataFrames, donc un état bricolé à la main
# masquerait des bugs que seul un enchaînement réel révèle.
from pathlib import Path

import openpyxl
from pptx import Presentation

from agent.nodes.build import build_node
from agent.nodes.export import export_node
from agent.nodes.test import test_node as run_test_node
from agent.nodes.validate import validate_node
from agent.state import AgentState, AnalysisStatus
from agent.tools.data_loader import load_data


def _fully_inspected_state(csv_path, tmp_path, **overrides):
    db_path = str(tmp_path / "analytics.duckdb")
    meta = load_data(csv_path, db_path=db_path)
    state = AgentState(
        query="q", data_path=csv_path, data_metadata=meta,
        business_question="Pourquoi la rétention baisse-t-elle ?",
        metric_definition="Rétention hebdomadaire",
        **overrides,
    )
    state_dict = state.model_dump()
    state_dict.update(build_node(state))
    state = AgentState(**state_dict)
    state_dict.update(run_test_node(state))
    state = AgentState(**state_dict)
    state_dict.update(validate_node(state))
    return AgentState(**state_dict)


def test_export_node_generates_valid_excel_and_pptx(monkeypatch, sample_retention_csv, tmp_path):
    monkeypatch.chdir(tmp_path)  # export_node écrit dans ./outputs (chemin relatif, non paramétrable)
    state = _fully_inspected_state(sample_retention_csv, tmp_path, recommendations=[
        {"title": "Investiguer le mobile", "description": "Détail.",
         "impact": "élevé", "feasibility": "moyenne", "timeline": "court terme"},
    ])

    result = export_node(state)

    assert result["status"] == AnalysisStatus.COMPLETED
    assert result.get("errors", []) == []
    assert Path(result["excel_path"]).exists()
    assert Path(result["presentation_path"]).exists()

    wb = openpyxl.load_workbook(result["excel_path"])
    assert "Validation" in wb.sheetnames
    assert any(s != "Sheet" for s in wb.sheetnames)

    prs = Presentation(result["presentation_path"])
    assert len(prs.slides) == 6  # PresentationGenerator produit toujours 6 diapositives

    assert any("Export" in line for line in result["audit_trail"])


def test_export_node_reports_consistency_check_result(monkeypatch, sample_retention_csv, tmp_path):
    monkeypatch.chdir(tmp_path)
    state = _fully_inspected_state(sample_retention_csv, tmp_path)

    result = export_node(state)

    # Les chiffres du PPTX proviennent des mêmes données que l'Excel (mêmes
    # requêtes ré-exécutées) : la vérification de cohérence doit passer.
    assert any("cohérence Excel/PPTX OK" in line for line in result["audit_trail"])


def test_export_node_handles_state_without_optional_fields(monkeypatch, tmp_path):
    # Ni weekly_retention, ni driver_analysis, ni recommendations -- ne doit
    # pas planter, juste produire des livrables presque vides.
    monkeypatch.chdir(tmp_path)
    state = AgentState(
        query="q", data_path="data.csv", data_metadata={"db_path": None},
        business_question="q", metric_definition="m",
    )

    result = export_node(state)

    assert result["status"] == AnalysisStatus.COMPLETED
    assert Path(result["excel_path"]).exists()
    assert Path(result["presentation_path"]).exists()

    prs = Presentation(result["presentation_path"])
    findings_slide_text = prs.slides[1].shapes[1].text_frame.text
    assert "Aucune observation disponible" in findings_slide_text


def test_export_node_records_error_without_blocking_pipeline(monkeypatch, tmp_path):
    # Une requête "query" invalide fait planter fetch_dataframe() --
    # l'export échoue mais l'analyse reste COMPLETED (l'export est un bonus,
    # pas le livrable principal -- voir la docstring d'export_node).
    monkeypatch.chdir(tmp_path)
    state = AgentState(
        query="q", data_path="data.csv",
        data_metadata={"db_path": str(tmp_path / "inexistant.duckdb")},
        business_question="q", metric_definition="m",
        weekly_retention={"data": "...", "label": "Test", "query": "SELECT * FROM table_qui_nexiste_pas"},
    )

    result = export_node(state)

    assert result["status"] == AnalysisStatus.COMPLETED
    assert len(result["errors"]) == 1
    assert "Export Excel/PPTX échoué" in result["errors"][0]
    assert any("échec" in line for line in result["audit_trail"])
