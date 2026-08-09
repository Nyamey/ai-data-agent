# tests/test_consistency_check.py — Cohérence des chiffres Excel <-> PowerPoint
#
# verify_consistency() compare des nombres à la granularité d'un "run" pptx
# entier (pas d'extraction de nombres au sein d'une phrase) : ces tests
# construisent donc des présentations avec un nombre isolé dans son propre
# passage de texte, pour cibler précisément ce mécanisme de comparaison.
import pandas as pd
from pptx import Presentation
from pptx.util import Inches

from agent.output.consistency_check import verify_consistency
from agent.output.excel_generator import ExcelGenerator


def _build_excel(tmp_path, values):
    gen = ExcelGenerator()
    gen.add_weekly_retention(pd.DataFrame({"semaine": range(len(values)), "valeur": values}))
    path = str(tmp_path / "rapport.xlsx")
    gen.save(path)
    return path


def _build_pptx_with_number(tmp_path, number):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text_frame.paragraphs[0].text = str(number)
    path = str(tmp_path / "presentation.pptx")
    prs.save(path)
    return path


def test_consistency_passes_when_pptx_number_exists_in_excel(tmp_path):
    excel_path = _build_excel(tmp_path, [10, 20, 30])
    pptx_path = _build_pptx_with_number(tmp_path, 20)

    result = verify_consistency(excel_path, pptx_path)
    assert result["passed"] is True
    assert result["excel_count"] > 0


def test_consistency_flags_number_absent_from_excel(tmp_path):
    excel_path = _build_excel(tmp_path, [10, 20, 30])
    pptx_path = _build_pptx_with_number(tmp_path, 999999)

    result = verify_consistency(excel_path, pptx_path)
    assert result["passed"] is False
    assert 999999.0 in result["inconsistencies"]
